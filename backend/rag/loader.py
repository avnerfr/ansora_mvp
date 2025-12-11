from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from typing import List
from langchain_core.documents import Document
import os

# Optional imports - gracefully handle if not available
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import docx  # type: ignore[import]
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False


def load_document(file_path: str, file_type: str) -> List[Document]:
    """Load a document and return LangChain Document objects."""
    documents = []
    
    try:
        if file_type in ["application/pdf", "pdf"]:
            loader = PyPDFLoader(file_path)
            documents = loader.load()
        elif file_type in ["text/plain", "txt"]:
            # Simple text file loading
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            except Exception as e:
                text = f"Error reading text file: {str(e)}"
            documents = [Document(page_content=text, metadata={"source": file_path})]
        elif file_type in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
            "docx",
            "doc",
        ]:
            # Word documents (.docx / .doc)
            if DOCX_AVAILABLE:
                try:
                    doc = docx.Document(file_path)
                    text_parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                    content = "\n".join(text_parts) if text_parts else "No text content found in document."
                    documents = [Document(page_content=content, metadata={"source": file_path})]
                except Exception as e:
                    documents = [Document(
                        page_content=f"Error loading Word document: {str(e)}",
                        metadata={"source": file_path, "error": True}
                    )]
            else:
                documents = [Document(
                    page_content=f"Word document: {os.path.basename(file_path)}. python-docx not available.",
                    metadata={"source": file_path}
                )]
        elif file_type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
            "ppt"
        ]:
            # Use lightweight python-pptx instead of heavy unstructured library
            if PPTX_AVAILABLE:
                try:
                    prs = Presentation(file_path)
                    text_parts = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                        if slide_text:
                            text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
                    content = "\n\n".join(text_parts) if text_parts else "No text content found in presentation."
                    documents = [Document(page_content=content, metadata={"source": file_path})]
                except Exception as e:
                    documents = [Document(
                        page_content=f"Error loading PowerPoint file: {str(e)}",
                        metadata={"source": file_path, "error": True}
                    )]
            else:
                documents = [Document(
                    page_content=f"PowerPoint file: {os.path.basename(file_path)}. python-pptx not available.",
                    metadata={"source": file_path}
                )]
        elif file_type in ["image/png", "image/jpeg", "image/jpg", "image/gif"]:
            # Basic OCR for images - optional feature
            if OCR_AVAILABLE:
                try:
                    image = Image.open(file_path)
                    text = pytesseract.image_to_string(image)
                    documents = [Document(page_content=text, metadata={"source": file_path})]
                except Exception as e:
                    documents = [Document(
                        page_content=f"Image file: {os.path.basename(file_path)}. OCR extraction failed: {str(e)}",
                        metadata={"source": file_path}
                    )]
            else:
                documents = [Document(
                    page_content=f"Image file: {os.path.basename(file_path)}. OCR not available (pytesseract not installed).",
                    metadata={"source": file_path}
                )]
        else:
            # For other file types, create a placeholder document
            documents = [Document(
                page_content=f"File content not extracted: {os.path.basename(file_path)}",
                metadata={"source": file_path}
            )]
    except Exception as e:
        # Return a document with error info
        documents = [Document(
            page_content=f"Error loading document: {str(e)}",
            metadata={"source": file_path, "error": True}
        )]
    
    return documents


def chunk_documents(documents: List[Document], chunk_size: int = 1000, chunk_overlap: int = 0) -> List[Document]:
    """
    Split documents into chunks.
    
    Notes:
    - We use chunk_overlap=0 to avoid duplicated/overlapping text segments
      showing up multiple times in the RAG prompt.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )
    chunks = text_splitter.split_documents(documents)
    return chunks

